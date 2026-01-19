"""
Bheem Workspace - Unified Presentation Service
==============================================
Handles presentation operations for both Internal (ERP) and External (SaaS) modes.
Integrates with OnlyOffice Document Server for full PowerPoint-compatible editing.

Storage Architecture:
- Internal mode: S3 at internal/{company_id}/presentations/
- External mode: S3 at external/{tenant_id}/presentations/
- Nextcloud sync: Optional WebDAV sync for user access

Features:
- PPTX file-based storage (not JSON blobs)
- OnlyOffice real-time collaboration
- Version control
- ERP entity linking (internal mode)
- Nextcloud integration
"""

import logging
import hashlib
import uuid
import jwt
import time
from datetime import datetime, timedelta
from enum import Enum
from io import BytesIO
from typing import Optional, Dict, Any, List
from uuid import UUID

import httpx
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession

from core.config import settings
from services.docs_storage_service import DocsStorageService, get_docs_storage_service
from services.nextcloud_service import NextcloudService

logger = logging.getLogger(__name__)


class PresentationMode(str, Enum):
    """Presentation storage mode"""
    INTERNAL = "internal"  # ERP mode - uses company_id, entity linking
    EXTERNAL = "external"  # SaaS mode - uses tenant_id, simple sharing


class PresentationService:
    """
    Unified presentation service for dual-mode architecture.

    Internal Mode (ERP):
    - Storage: S3 at internal/{company_id}/presentations/
    - Database: workspace.presentations with storage_mode='internal'
    - Features: ERP entity linking, versioning, audit
    - Users: ERP employees synced from HR

    External Mode (SaaS):
    - Storage: S3 at external/{tenant_id}/presentations/
    - Database: workspace.presentations with storage_mode='external'
    - Features: Sharing, collaboration, public links
    - Users: Self-registered SaaS customers
    """

    def __init__(self, db: AsyncSession):
        self.db = db
        self.storage = get_docs_storage_service()
        self.nextcloud = NextcloudService() if settings.PROVISIONING_SYNC_NEXTCLOUD else None

    # =============================================
    # Core Operations
    # =============================================

    async def create_presentation(
        self,
        title: str,
        mode: PresentationMode,
        tenant_id: UUID,
        user_id: UUID,
        company_id: Optional[UUID] = None,
        description: Optional[str] = None,
        folder_id: Optional[UUID] = None,
        template_id: Optional[str] = None,
        linked_entity_type: Optional[str] = None,
        linked_entity_id: Optional[UUID] = None,
    ) -> Dict[str, Any]:
        """
        Create a new presentation as PPTX file.

        Args:
            title: Presentation title
            mode: Internal (ERP) or External (SaaS)
            tenant_id: Workspace tenant ID
            user_id: Creator's user ID
            company_id: ERP company ID (required for internal mode)
            description: Optional description
            folder_id: Optional folder for organization
            template_id: Optional template to use
            linked_entity_type: ERP entity type (internal mode)
            linked_entity_id: ERP entity ID (internal mode)

        Returns:
            Presentation metadata with edit URL
        """
        presentation_id = uuid.uuid4()
        document_key = self._generate_document_key(presentation_id)

        # Generate PPTX content
        if template_id:
            pptx_bytes = await self._get_template_content(template_id)
        else:
            pptx_bytes = self._create_empty_pptx(title)

        # Calculate checksum
        checksum = hashlib.sha256(pptx_bytes).hexdigest()
        file_size = len(pptx_bytes)

        # Determine storage path based on mode
        if mode == PresentationMode.INTERNAL:
            if not company_id:
                raise ValueError("company_id is required for internal mode")
            storage_path = f"internal/{company_id}/presentations/{presentation_id}.pptx"
        else:
            storage_path = f"external/{tenant_id}/presentations/{presentation_id}.pptx"

        # Upload to S3
        try:
            await self.storage.upload_file(
                file=BytesIO(pptx_bytes),
                filename=f"{presentation_id}.pptx",
                company_id=company_id if mode == PresentationMode.INTERNAL else None,
                tenant_id=tenant_id if mode == PresentationMode.EXTERNAL else None,
                folder_path="presentations",
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            logger.info(f"Uploaded presentation to S3: {storage_path}")
        except Exception as e:
            logger.error(f"Failed to upload presentation to S3: {e}")
            raise

        # Default theme
        default_theme = {
            "font_heading": "Arial",
            "font_body": "Arial",
            "color_primary": "#1a73e8",
            "color_secondary": "#34a853",
            "color_background": "#ffffff"
        }

        # Insert into database
        query = text("""
            INSERT INTO workspace.presentations (
                id, tenant_id, title, description, folder_id, theme,
                storage_path, storage_bucket, file_size, checksum,
                version, storage_mode, document_key,
                linked_entity_type, linked_entity_id,
                created_by, created_at, updated_at
            ) VALUES (
                :id, :tenant_id, :title, :description, :folder_id, :theme,
                :storage_path, :storage_bucket, :file_size, :checksum,
                1, :storage_mode, :document_key,
                :linked_entity_type, :linked_entity_id,
                :user_id, NOW(), NOW()
            )
            RETURNING id, title, storage_path, version, document_key, created_at
        """)

        result = await self.db.execute(query, {
            "id": presentation_id,
            "tenant_id": tenant_id,
            "title": title,
            "description": description,
            "folder_id": folder_id,
            "theme": str(default_theme),
            "storage_path": storage_path,
            "storage_bucket": self.storage.bucket,
            "file_size": file_size,
            "checksum": checksum,
            "storage_mode": mode.value,
            "document_key": document_key,
            "linked_entity_type": linked_entity_type,
            "linked_entity_id": linked_entity_id,
            "user_id": user_id,
        })
        await self.db.commit()

        row = result.fetchone()

        # Create initial version record
        await self._create_version_record(
            presentation_id=presentation_id,
            version_number=1,
            storage_path=storage_path,
            file_size=file_size,
            checksum=checksum,
            user_id=user_id,
            comment="Initial version"
        )

        # Create default title slide in slides table
        slide_id = uuid.uuid4()
        await self.db.execute(text("""
            INSERT INTO workspace.slides
            (id, presentation_id, slide_index, layout, content, created_at, updated_at)
            VALUES (
                :id, :presentation_id, 0, 'title',
                :content,
                NOW(), NOW()
            )
        """), {
            "id": slide_id,
            "presentation_id": presentation_id,
            "content": '{"title": {"text": "' + title + '"}, "subtitle": {"text": "Click to add subtitle"}}'
        })
        await self.db.commit()

        return {
            "id": str(presentation_id),
            "title": title,
            "description": description,
            "storage_path": storage_path,
            "storage_mode": mode.value,
            "version": 1,
            "document_key": document_key,
            "file_size": file_size,
            "created_at": row.created_at.isoformat() if row.created_at else None,
            "edit_url": f"/slides/{presentation_id}/edit",
        }

    async def get_presentation(
        self,
        presentation_id: UUID,
        tenant_id: UUID,
        user_id: UUID,
    ) -> Optional[Dict[str, Any]]:
        """Get presentation metadata."""
        query = text("""
            SELECT
                p.id, p.title, p.description, p.folder_id, p.theme,
                p.storage_path, p.storage_bucket, p.file_size, p.checksum,
                p.nextcloud_path, p.version, p.storage_mode, p.document_key,
                p.linked_entity_type, p.linked_entity_id,
                p.is_starred, p.is_deleted,
                p.created_by, p.created_at, p.updated_at,
                p.last_edited_by, p.last_edited_at,
                tu.email as creator_email,
                tu.name as creator_name
            FROM workspace.presentations p
            LEFT JOIN workspace.tenant_users tu ON p.created_by = tu.id
            WHERE p.id = :presentation_id
            AND p.tenant_id = :tenant_id
            AND p.is_deleted = FALSE
        """)

        result = await self.db.execute(query, {
            "presentation_id": presentation_id,
            "tenant_id": tenant_id,
        })
        row = result.fetchone()

        if not row:
            return None

        return {
            "id": str(row.id),
            "title": row.title,
            "description": row.description,
            "folder_id": str(row.folder_id) if row.folder_id else None,
            "theme": row.theme,
            "storage_path": row.storage_path,
            "storage_bucket": row.storage_bucket,
            "file_size": row.file_size,
            "checksum": row.checksum,
            "nextcloud_path": row.nextcloud_path,
            "version": row.version or 1,
            "storage_mode": row.storage_mode,
            "document_key": row.document_key,
            "linked_entity_type": row.linked_entity_type,
            "linked_entity_id": str(row.linked_entity_id) if row.linked_entity_id else None,
            "is_starred": row.is_starred,
            "created_by": str(row.created_by),
            "creator_email": row.creator_email,
            "creator_name": row.creator_name,
            "created_at": row.created_at.isoformat() if row.created_at else None,
            "updated_at": row.updated_at.isoformat() if row.updated_at else None,
            "last_edited_by": str(row.last_edited_by) if row.last_edited_by else None,
            "last_edited_at": row.last_edited_at.isoformat() if row.last_edited_at else None,
        }

    # =============================================
    # OnlyOffice Integration
    # =============================================

    async def get_editor_config(
        self,
        presentation_id: UUID,
        tenant_id: UUID,
        user_id: UUID,
        user_name: str,
        user_email: str,
        mode: str = "edit",  # edit, view, review
    ) -> Dict[str, Any]:
        """
        Get OnlyOffice Document Editor configuration.

        This generates the config object that the OnlyOffice JS API needs
        to initialize the editor.

        Args:
            presentation_id: Presentation UUID
            tenant_id: Tenant UUID for access validation
            user_id: User UUID
            user_name: Display name for collaboration
            user_email: User email
            mode: Edit mode (edit, view, review)

        Returns:
            OnlyOffice editor configuration
        """
        # Get presentation metadata
        presentation = await self.get_presentation(presentation_id, tenant_id, user_id)
        if not presentation:
            raise ValueError("Presentation not found")

        # Check if storage_path exists, if not create the PPTX file
        storage_path = presentation.get("storage_path")
        if not storage_path:
            # Create empty PPTX file for legacy presentations
            storage_path = await self._create_pptx_for_legacy_presentation(
                presentation_id=presentation_id,
                tenant_id=tenant_id,
                user_id=user_id,
                title=presentation.get("title", "Untitled"),
            )
            presentation["storage_path"] = storage_path

        # Generate document key (unique per version for cache invalidation)
        version = presentation.get('version') or 1
        document_key = f"{presentation_id}-v{version}-{int(time.time())}"

        # Generate access token for document proxy endpoint
        access_token = jwt.encode(
            {
                "presentation_id": str(presentation_id),
                "exp": datetime.utcnow() + timedelta(hours=24),
            },
            settings.ONLYOFFICE_JWT_SECRET,
            algorithm="HS256"
        )

        # Ensure token is a string (PyJWT 2.0+ returns string, older versions returned bytes)
        if isinstance(access_token, bytes):
            access_token = access_token.decode('utf-8')

        # URL-encode the token for safety (even though base64url should be URL-safe)
        from urllib.parse import quote
        encoded_token = quote(access_token, safe='')

        # Use proxy URL instead of S3 presigned URL (OnlyOffice can't access S3 directly)
        # The proxy endpoint streams the file from S3 through our backend
        download_url = f"{settings.WORKSPACE_URL}/api/v1/slides/{presentation_id}/content?token={encoded_token}"

        logger.info(f"Generated download URL for presentation {presentation_id}: {download_url[:100]}...")

        # Update document key in database
        await self.db.execute(
            text("UPDATE workspace.presentations SET document_key = :key WHERE id = :id"),
            {"key": document_key, "id": presentation_id}
        )
        await self.db.commit()

        # Callback URL for OnlyOffice to notify us of changes
        callback_url = f"{settings.ONLYOFFICE_CALLBACK_URL.replace('/sheets', '/slides')}/{presentation_id}/onlyoffice-callback"

        # Build editor config
        config = {
            "document": {
                "fileType": "pptx",
                "key": document_key,
                "title": f"{presentation['title']}.pptx",
                "url": download_url,
                "permissions": {
                    "comment": True,
                    "download": True,
                    "edit": mode == "edit",
                    "print": True,
                    "review": mode in ["edit", "review"],
                },
            },
            "documentType": "slide",  # presentation type
            "editorConfig": {
                "callbackUrl": callback_url,
                "lang": "en",
                "mode": mode,
                "user": {
                    "id": str(user_id),
                    "name": user_name,
                },
                "customization": {
                    "autosave": True,
                    "chat": True,
                    "comments": True,
                    "compactHeader": False,
                    "compactToolbar": False,
                    "feedback": False,
                    "forcesave": True,
                    "help": False,  # Hide OnlyOffice help
                    "hideRightMenu": False,
                    "about": False,  # Hide OnlyOffice about dialog
                    "logo": {
                        "image": "https://workspace.bheem.cloud/bheem-logo.png",
                        "imageEmbedded": "https://workspace.bheem.cloud/bheem-logo-small.png",
                        "url": "https://workspace.bheem.cloud",
                    },
                    "loaderLogo": "https://workspace.bheem.cloud/bheem-logo.png",
                    "loaderName": "Bheem Slides",
                    "customer": {
                        "address": "Bheem Workspace",
                        "logo": "https://workspace.bheem.cloud/bheem-logo.png",
                        "logoEmbedded": "https://workspace.bheem.cloud/bheem-logo-small.png",
                        "mail": "support@bheem.cloud",
                        "name": "Bheem Slides",
                        "www": "https://workspace.bheem.cloud",
                    },
                    "toolbarNoTabs": False,
                    "zoom": 100,
                },
            },
            "height": "100%",
            "width": "100%",
            "type": "desktop",
        }

        # Add JWT token if enabled
        if settings.ONLYOFFICE_JWT_ENABLED and settings.ONLYOFFICE_JWT_SECRET:
            token = self._generate_onlyoffice_jwt(config)
            config["token"] = token

        # Record edit session
        await self._record_edit_session(
            presentation_id=presentation_id,
            user_id=user_id,
            document_key=document_key,
        )

        return {
            "config": config,
            "documentServerUrl": settings.ONLYOFFICE_URL,
            "presentation": presentation,
        }

    async def handle_onlyoffice_callback(
        self,
        presentation_id: UUID,
        callback_data: Dict[str, Any],
    ) -> Dict[str, int]:
        """
        Handle OnlyOffice Document Server callback.

        OnlyOffice calls this endpoint when document state changes.

        Status codes:
        - 0: No document change
        - 1: Document being edited
        - 2: Document ready for saving (users closed editors)
        - 3: Document saving error
        - 4: Document closed with no changes
        - 6: Document being edited, force save requested
        - 7: Error force saving document
        """
        status = callback_data.get("status")
        document_key = callback_data.get("key")

        logger.info(f"OnlyOffice callback for presentation {presentation_id}: status={status}, key={document_key}")

        try:
            if status == 2 or status == 6:
                # Document ready for saving or force save
                document_url = callback_data.get("url")

                if not document_url:
                    logger.error("No document URL in callback")
                    return {"error": 1}

                # Download the edited document from OnlyOffice
                async with httpx.AsyncClient(verify=False) as client:
                    response = await client.get(document_url, timeout=60.0)
                    if response.status_code != 200:
                        logger.error(f"Failed to download from OnlyOffice: {response.status_code}")
                        return {"error": 1}
                    new_pptx_bytes = response.content

                # Get current presentation info
                result = await self.db.execute(
                    text("""
                        SELECT storage_path, storage_bucket, version, tenant_id, created_by
                        FROM workspace.presentations
                        WHERE id = :id
                    """),
                    {"id": presentation_id}
                )
                row = result.fetchone()

                if not row:
                    logger.error(f"Presentation not found: {presentation_id}")
                    return {"error": 1}

                storage_path = row.storage_path
                current_version = row.version or 1

                # Calculate new checksum
                new_checksum = hashlib.sha256(new_pptx_bytes).hexdigest()
                new_file_size = len(new_pptx_bytes)

                # Upload new version to S3 (overwrite)
                await self.storage.upload_file(
                    file=BytesIO(new_pptx_bytes),
                    filename=storage_path.split("/")[-1],
                    company_id=None,
                    tenant_id=None,
                    folder_path="/".join(storage_path.split("/")[:-1]),
                    content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

                # Increment version
                new_version = current_version + 1

                # Update database
                await self.db.execute(
                    text("""
                        UPDATE workspace.presentations
                        SET version = :version,
                            file_size = :file_size,
                            checksum = :checksum,
                            updated_at = NOW(),
                            last_edited_at = NOW()
                        WHERE id = :id
                    """),
                    {
                        "id": presentation_id,
                        "version": new_version,
                        "file_size": new_file_size,
                        "checksum": new_checksum,
                    }
                )

                # Create version record
                await self._create_version_record(
                    presentation_id=presentation_id,
                    version_number=new_version,
                    storage_path=storage_path,
                    file_size=new_file_size,
                    checksum=new_checksum,
                    user_id=row.created_by,
                    comment=f"Auto-saved version {new_version}"
                )

                await self.db.commit()

                logger.info(f"Saved presentation {presentation_id} version {new_version}")

            elif status == 4:
                # Document closed with no changes
                logger.info(f"Presentation {presentation_id} closed without changes")

            elif status == 1:
                # Document being edited - update session
                await self._update_edit_session(presentation_id, document_key)

            return {"error": 0}

        except Exception as e:
            logger.error(f"OnlyOffice callback error: {e}")
            await self.db.rollback()
            return {"error": 1}

    # =============================================
    # Helper Methods
    # =============================================

    def _create_empty_pptx(self, title: str) -> bytes:
        """Create empty PPTX presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            # If python-pptx is not installed, create a minimal PPTX
            logger.warning("python-pptx not installed, creating minimal PPTX")
            return self._create_minimal_pptx()

        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(title_slide_layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = "Click to add subtitle"

        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _create_minimal_pptx(self) -> bytes:
        """Create minimal PPTX without python-pptx."""
        # This is a minimal valid PPTX file structure
        import zipfile
        buffer = BytesIO()

        with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
            # [Content_Types].xml
            content_types = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
    <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
    <Default Extension="xml" ContentType="application/xml"/>
    <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
    <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
    <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
    <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
</Types>'''
            zf.writestr('[Content_Types].xml', content_types)

            # _rels/.rels
            rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>'''
            zf.writestr('_rels/.rels', rels)

            # ppt/presentation.xml
            presentation = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
    <p:sldIdLst><p:sldId id="256" r:id="rId2"/></p:sldIdLst>
    <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
</p:presentation>'''
            zf.writestr('ppt/presentation.xml', presentation)

            # ppt/_rels/presentation.xml.rels
            pres_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
    <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>'''
            zf.writestr('ppt/_rels/presentation.xml.rels', pres_rels)

            # ppt/slides/slide1.xml
            slide1 = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
    <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
</p:sld>'''
            zf.writestr('ppt/slides/slide1.xml', slide1)

            # ppt/slides/_rels/slide1.xml.rels
            slide1_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''
            zf.writestr('ppt/slides/_rels/slide1.xml.rels', slide1_rels)

            # ppt/slideMasters/slideMaster1.xml
            master = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
    <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
    <p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
</p:sldMaster>'''
            zf.writestr('ppt/slideMasters/slideMaster1.xml', master)

            # ppt/slideMasters/_rels/slideMaster1.xml.rels
            master_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''
            zf.writestr('ppt/slideMasters/_rels/slideMaster1.xml.rels', master_rels)

            # ppt/slideLayouts/slideLayout1.xml
            layout = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" type="title">
    <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
</p:sldLayout>'''
            zf.writestr('ppt/slideLayouts/slideLayout1.xml', layout)

            # ppt/slideLayouts/_rels/slideLayout1.xml.rels
            layout_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>'''
            zf.writestr('ppt/slideLayouts/_rels/slideLayout1.xml.rels', layout_rels)

        return buffer.getvalue()

    async def _get_template_content(self, template_id: str) -> bytes:
        """Get template content from database or storage."""
        result = await self.db.execute(
            text("""
                SELECT content FROM workspace.productivity_templates
                WHERE id = :id AND template_type = 'presentation'
            """),
            {"id": template_id}
        )
        row = result.fetchone()

        if row and row.content:
            if isinstance(row.content, dict) and "storage_path" in row.content:
                file_content, _ = await self.storage.download_file(row.content["storage_path"])
                return file_content.read()

        return self._create_empty_pptx("Untitled presentation")

    def _generate_document_key(self, presentation_id: UUID) -> str:
        """Generate unique document key for OnlyOffice."""
        timestamp = int(time.time() * 1000)
        return f"{presentation_id}-{timestamp}"

    def _generate_onlyoffice_jwt(self, config: Dict[str, Any]) -> str:
        """Generate JWT token for OnlyOffice."""
        payload = config.copy()
        payload["exp"] = datetime.utcnow() + timedelta(hours=24)

        return jwt.encode(
            payload,
            settings.ONLYOFFICE_JWT_SECRET,
            algorithm="HS256"
        )

    async def _create_pptx_for_legacy_presentation(
        self,
        presentation_id: UUID,
        tenant_id: UUID,
        user_id: UUID,
        title: str,
    ) -> str:
        """
        Create PPTX file for legacy presentations that don't have storage_path.
        """
        # Create workbook
        pptx_bytes = self._create_empty_pptx(title)
        checksum = hashlib.sha256(pptx_bytes).hexdigest()
        file_size = len(pptx_bytes)

        # Storage path
        storage_path = f"external/{tenant_id}/presentations/{presentation_id}.pptx"

        # Upload to S3
        await self.storage.upload_file(
            file=BytesIO(pptx_bytes),
            filename=f"{presentation_id}.pptx",
            tenant_id=tenant_id,
            folder_path="presentations",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        # Update database
        await self.db.execute(
            text("""
                UPDATE workspace.presentations
                SET storage_path = :storage_path,
                    storage_bucket = :storage_bucket,
                    file_size = :file_size,
                    checksum = :checksum,
                    version = COALESCE(version, 1),
                    storage_mode = COALESCE(storage_mode, 'external')
                WHERE id = :id
            """),
            {
                "id": presentation_id,
                "storage_path": storage_path,
                "storage_bucket": self.storage.bucket,
                "file_size": file_size,
                "checksum": checksum,
            }
        )
        await self.db.commit()

        logger.info(f"Created PPTX for legacy presentation {presentation_id}")
        return storage_path

    async def _create_version_record(
        self,
        presentation_id: UUID,
        version_number: int,
        storage_path: str,
        file_size: int,
        checksum: str,
        user_id: UUID,
        comment: str = None,
    ):
        """Create version history record."""
        await self.db.execute(
            text("""
                INSERT INTO workspace.presentation_versions
                (presentation_id, version_number, storage_path, file_size, checksum, created_by, comment)
                VALUES (:presentation_id, :version_number, :storage_path, :file_size, :checksum, :created_by, :comment)
                ON CONFLICT (presentation_id, version_number) DO NOTHING
            """),
            {
                "presentation_id": presentation_id,
                "version_number": version_number,
                "storage_path": storage_path,
                "file_size": file_size,
                "checksum": checksum,
                "created_by": user_id,
                "comment": comment,
            }
        )

    async def _record_edit_session(
        self,
        presentation_id: UUID,
        user_id: UUID,
        document_key: str,
    ):
        """Record active editing session."""
        await self.db.execute(
            text("""
                INSERT INTO workspace.presentation_edit_sessions
                (presentation_id, user_id, session_key, started_at, last_activity, is_active)
                VALUES (:presentation_id, :user_id, :session_key, NOW(), NOW(), TRUE)
                ON CONFLICT DO NOTHING
            """),
            {
                "presentation_id": presentation_id,
                "user_id": user_id,
                "session_key": document_key,
            }
        )

    async def _update_edit_session(
        self,
        presentation_id: UUID,
        document_key: str,
    ):
        """Update edit session last activity."""
        await self.db.execute(
            text("""
                UPDATE workspace.presentation_edit_sessions
                SET last_activity = NOW()
                WHERE presentation_id = :presentation_id
                AND session_key = :session_key
                AND is_active = TRUE
            """),
            {
                "presentation_id": presentation_id,
                "session_key": document_key,
            }
        )


# Factory function
def get_presentation_service(db: AsyncSession) -> PresentationService:
    """Get presentation service instance."""
    return PresentationService(db)
